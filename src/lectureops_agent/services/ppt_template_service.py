from __future__ import annotations

import hashlib
import io
import os
import re
import zipfile
from datetime import datetime, timezone
from typing import Any, Protocol
from uuid import uuid4

from pptx import Presentation

from lectureops_agent.models.schemas import PPTTemplateLayout, PPTTemplateMetadata


PPT_TEMPLATE_SEMANTIC_TYPES = (
    "cover",
    "objectives",
    "lesson",
    "practice",
    "assessment",
    "ncs_coverage",
    "sources",
)
SOURCE_SLIDE_LAYOUT_OFFSET = 10_000
DEFAULT_MAX_ZIP_ENTRIES = 5000
DEFAULT_MAX_UNCOMPRESSED_BYTES = 250 * 1024 * 1024
_SAFE_FILENAME = re.compile(r"[^0-9A-Za-z가-힣._ -]+")
_TITLE_TYPES = {"TITLE", "CENTER_TITLE"}
_BODY_TYPES = {"BODY", "OBJECT", "SUBTITLE", "VERTICAL_BODY", "VERTICAL_OBJECT"}
_SOURCE_COVER_EXCLUSION_TERMS = (
    "contents",
    "table of contents",
    "목차",
    "free fonts",
    "template uses",
    "presentation template is free",
    "happy designing",
    "license",
    "라이선스",
    "사용 안내",
    "사용방법",
)
_SOURCE_AUTHORING_EXCLUSION_TERMS = (
    "resource page",
    "credits",
    "free fonts",
    "template uses",
    "presentation template is free",
    "slides carnival",
    "happy designing",
    "라이선스",
    "사용 안내",
    "사용방법",
)
_SOURCE_ROLE_LABELS = {
    "cover": "표지",
    "objectives": "학습목표·개요",
    "lesson": "교안 본문",
    "practice": "실습·절차",
    "assessment": "평가",
    "ncs_coverage": "NCS·핵심정리",
    "sources": "출처·마무리",
}


class PPTTemplateStore(Protocol):
    def readiness(self) -> dict[str, Any]:
        ...

    def save(self, metadata: PPTTemplateMetadata, content: bytes) -> PPTTemplateMetadata:
        ...

    def get(self, project_id: str) -> PPTTemplateMetadata | None:
        ...

    def load_content(self, project_id: str) -> bytes:
        ...

    def update_mapping(
        self,
        project_id: str,
        layout_mapping: dict[str, int],
    ) -> PPTTemplateMetadata | None:
        ...

    def delete(self, project_id: str) -> bool:
        ...


class InMemoryPPTTemplateStore:
    def __init__(self) -> None:
        self._metadata_by_project: dict[str, PPTTemplateMetadata] = {}
        self._content_by_project: dict[str, bytes] = {}

    def readiness(self) -> dict[str, Any]:
        return {"ready": True, "store": type(self).__name__}

    def save(self, metadata: PPTTemplateMetadata, content: bytes) -> PPTTemplateMetadata:
        self._metadata_by_project[metadata.project_id] = metadata
        self._content_by_project[metadata.project_id] = bytes(content)
        return metadata

    def get(self, project_id: str) -> PPTTemplateMetadata | None:
        return self._metadata_by_project.get(project_id)

    def load_content(self, project_id: str) -> bytes:
        try:
            return self._content_by_project[project_id]
        except KeyError as exc:
            raise KeyError("PPT template not found") from exc

    def update_mapping(
        self,
        project_id: str,
        layout_mapping: dict[str, int],
    ) -> PPTTemplateMetadata | None:
        metadata = self.get(project_id)
        if metadata is None:
            return None
        updated = metadata.model_copy(
            update={
                "layout_mapping": dict(layout_mapping),
                "updated_at": datetime.now(timezone.utc),
            }
        )
        self._metadata_by_project[project_id] = updated
        return updated

    def delete(self, project_id: str) -> bool:
        existed = project_id in self._metadata_by_project
        self._metadata_by_project.pop(project_id, None)
        self._content_by_project.pop(project_id, None)
        return existed


class SupabasePPTTemplateStore:
    def __init__(
        self,
        *,
        client: Any,
        table_name: str = "lessonpack_ppt_templates",
        bucket_name: str = "lessonpack-ppt-templates",
    ) -> None:
        self._client = client
        self.table_name = table_name
        self.bucket_name = bucket_name

    def readiness(self) -> dict[str, Any]:
        try:
            self._client.table(self.table_name).select("project_id,template_id").limit(1).execute()
        except Exception as exc:
            return {
                "ready": False,
                "store": type(self).__name__,
                "table": self.table_name,
                "bucket": self.bucket_name,
                "error_type": type(exc).__name__,
            }
        return {
            "ready": True,
            "store": type(self).__name__,
            "table": self.table_name,
            "bucket": self.bucket_name,
        }

    def save(self, metadata: PPTTemplateMetadata, content: bytes) -> PPTTemplateMetadata:
        existing_row = self._get_row(metadata.project_id)
        storage_path = f"{metadata.project_id}/{metadata.template_id}.pptx"
        bucket = self._client.storage.from_(self.bucket_name)
        bucket.upload(
            storage_path,
            content,
            file_options={
                "content-type": (
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                ),
                "upsert": "false",
            },
        )
        row = _metadata_to_row(metadata, storage_path=storage_path)
        try:
            response = (
                self._client.table(self.table_name)
                .upsert(row, on_conflict="project_id")
                .execute()
            )
            _raise_for_supabase_error(response)
        except Exception:
            _best_effort_remove(bucket, storage_path)
            raise

        old_storage_path = str(existing_row.get("storage_path", "")) if existing_row else ""
        if old_storage_path and old_storage_path != storage_path:
            _best_effort_remove(bucket, old_storage_path)
        return metadata

    def get(self, project_id: str) -> PPTTemplateMetadata | None:
        row = self._get_row(project_id)
        return _metadata_from_row(row) if row else None

    def load_content(self, project_id: str) -> bytes:
        row = self._get_row(project_id)
        if row is None:
            raise KeyError("PPT template not found")
        content = self._client.storage.from_(self.bucket_name).download(row["storage_path"])
        if not isinstance(content, bytes):
            content = bytes(content)
        expected_hash = str(row["content_hash"])
        if hashlib.sha256(content).hexdigest() != expected_hash:
            raise RuntimeError("Stored PPT template hash mismatch")
        return content

    def update_mapping(
        self,
        project_id: str,
        layout_mapping: dict[str, int],
    ) -> PPTTemplateMetadata | None:
        if self._get_row(project_id) is None:
            return None
        updated_at = datetime.now(timezone.utc)
        response = (
            self._client.table(self.table_name)
            .update(
                {
                    "layout_mapping": dict(layout_mapping),
                    "updated_at": updated_at.isoformat(),
                }
            )
            .eq("project_id", project_id)
            .execute()
        )
        _raise_for_supabase_error(response)
        return self.get(project_id)

    def delete(self, project_id: str) -> bool:
        row = self._get_row(project_id)
        if row is None:
            return False
        self._client.storage.from_(self.bucket_name).remove([row["storage_path"]])
        response = (
            self._client.table(self.table_name)
            .delete()
            .eq("project_id", project_id)
            .execute()
        )
        _raise_for_supabase_error(response)
        return True

    def _get_row(self, project_id: str) -> dict[str, Any] | None:
        response = (
            self._client.table(self.table_name)
            .select("*")
            .eq("project_id", project_id)
            .limit(1)
            .execute()
        )
        rows = _response_data(response)
        return rows[0] if rows else None


def create_ppt_template_store_for_vector_store(vector_store: Any) -> PPTTemplateStore:
    from lectureops_agent.services.vector_store import SupabaseVectorStore

    if isinstance(vector_store, SupabaseVectorStore):
        return SupabasePPTTemplateStore(
            client=vector_store.client,
            table_name=os.getenv(
                "LESSONPACK_PPT_TEMPLATE_TABLE",
                "lessonpack_ppt_templates",
            ),
            bucket_name=os.getenv(
                "LESSONPACK_PPT_TEMPLATE_BUCKET",
                "lessonpack-ppt-templates",
            ),
        )
    return InMemoryPPTTemplateStore()


def analyze_ppt_template(
    *,
    project_id: str,
    filename: str,
    content: bytes,
) -> PPTTemplateMetadata:
    if not filename.casefold().endswith(".pptx"):
        raise ValueError("PPT template must use the .pptx format.")
    warnings = _validate_office_zip(content)

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise ValueError("The uploaded file is not a readable PPTX presentation.") from exc

    master_layouts = [
        _layout_manifest(index, layout)
        for index, layout in enumerate(presentation.slide_layouts)
    ]
    if not master_layouts:
        raise ValueError("PPT template must contain at least one slide layout.")

    source_cover_index = reusable_source_cover_index(presentation)
    source_layouts = _source_slide_manifests(
        presentation,
        source_cover_index=source_cover_index,
    )
    layouts = [*source_layouts, *master_layouts]
    if source_layouts:
        warnings.append(
            f"원본 슬라이드 디자인 {len(source_layouts)}개를 유형별 생성 후보로 분석했습니다."
        )
    elif len(presentation.slides) > 0:
        warnings.append(
            "안전하게 재사용할 원본 슬라이드 후보가 없어 마스터 레이아웃만 사용합니다."
        )
    if not any(layout.supports_title and layout.supports_body for layout in master_layouts):
        warnings.append(
            "제목+본문 placeholder 레이아웃이 없어 일부 슬라이드는 텍스트 상자를 사용합니다."
        )

    now = datetime.now(timezone.utc)
    return PPTTemplateMetadata(
        template_id=str(uuid4()),
        project_id=project_id,
        original_filename=_sanitize_filename(filename),
        content_hash=hashlib.sha256(content).hexdigest(),
        file_size_bytes=len(content),
        source_slide_count=len(presentation.slides),
        slide_width=presentation.slide_width,
        slide_height=presentation.slide_height,
        layouts=layouts,
        layout_mapping=_automatic_layout_mapping(layouts),
        warnings=warnings,
        created_at=now,
        updated_at=now,
    )


def reusable_source_cover_index(presentation: Presentation) -> int | None:
    """Return a safe source slide index whose visual design can become the cover."""
    for index, slide in enumerate(list(presentation.slides)[:3]):
        text = " ".join(
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ).casefold()
        if not text:
            continue
        if any(term in text for term in _SOURCE_COVER_EXCLUSION_TERMS):
            continue
        if any(getattr(rel, "is_external", False) for rel in slide.part.rels.values()):
            continue
        text_shape_count = sum(
            1
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        )
        if text_shape_count > 8:
            continue
        return index
    return None


def source_slide_index_for_layout(layout_index: int | None) -> int | None:
    if layout_index is None or layout_index < SOURCE_SLIDE_LAYOUT_OFFSET:
        return None
    return layout_index - SOURCE_SLIDE_LAYOUT_OFFSET


def resolve_template_mapping_for_export(
    presentation: Presentation,
    stored_mapping: dict[str, int] | None,
) -> dict[str, int]:
    mapping = dict(stored_mapping or {})
    if any(
        source_slide_index_for_layout(value) is not None
        for value in mapping.values()
    ):
        return mapping

    source_cover_index = reusable_source_cover_index(presentation)
    source_layouts = _source_slide_manifests(
        presentation,
        source_cover_index=source_cover_index,
    )
    if not source_layouts:
        return mapping
    master_layouts = [
        _layout_manifest(index, layout)
        for index, layout in enumerate(presentation.slide_layouts)
    ]
    return _automatic_layout_mapping([*source_layouts, *master_layouts])


def validate_layout_mapping(
    metadata: PPTTemplateMetadata,
    requested_mapping: dict[str, int],
) -> dict[str, int]:
    unknown_types = sorted(set(requested_mapping) - set(PPT_TEMPLATE_SEMANTIC_TYPES))
    if unknown_types:
        raise ValueError(f"Unsupported PPT semantic slide types: {', '.join(unknown_types)}")
    available_indices = {layout.layout_index for layout in metadata.layouts}
    invalid_indices = sorted(set(requested_mapping.values()) - available_indices)
    if invalid_indices:
        raise ValueError(
            "PPT layout indices do not exist in the uploaded template: "
            + ", ".join(str(index) for index in invalid_indices)
        )
    return {**metadata.layout_mapping, **requested_mapping}


def _validate_office_zip(content: bytes) -> list[str]:
    if not content.startswith(b"PK"):
        raise ValueError("The uploaded file is not an Office Open XML presentation.")
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            entries = archive.infolist()
            if len(entries) > DEFAULT_MAX_ZIP_ENTRIES:
                raise ValueError("PPTX archive contains too many entries.")
            uncompressed_size = sum(entry.file_size for entry in entries)
            if uncompressed_size > DEFAULT_MAX_UNCOMPRESSED_BYTES:
                raise ValueError("PPTX archive is too large after decompression.")
            names = {entry.filename.casefold() for entry in entries}
            if not any(name.startswith("ppt/presentation") for name in names):
                raise ValueError("PPTX archive does not contain a presentation document.")
            if any(name.endswith("vbaproject.bin") for name in names):
                raise ValueError("Macro-enabled PowerPoint templates are not supported.")
            has_external_relationship = any(
                entry.filename.casefold().endswith(".rels")
                and b'TargetMode="External"' in archive.read(entry)
                for entry in entries
            )
    except zipfile.BadZipFile as exc:
        raise ValueError("The uploaded PPTX archive is damaged.") from exc
    return (
        ["외부 파일·URL 관계가 감지되었습니다. 생성 슬라이드에는 해당 관계를 복사하지 않습니다."]
        if has_external_relationship
        else []
    )


def _layout_manifest(index: int, layout: Any) -> PPTTemplateLayout:
    placeholder_types: list[str] = []
    for placeholder in layout.placeholders:
        value = placeholder.placeholder_format.type
        placeholder_types.append(getattr(value, "name", str(value)))
    type_set = set(placeholder_types)
    body_count = sum(value in _BODY_TYPES for value in placeholder_types)
    return PPTTemplateLayout(
        layout_index=index,
        name=(layout.name or f"Layout {index + 1}").strip(),
        placeholder_count=len(placeholder_types),
        placeholder_types=placeholder_types,
        supports_title=bool(type_set & _TITLE_TYPES),
        supports_body=body_count > 0,
        body_placeholder_count=body_count,
    )


def _source_slide_manifests(
    presentation: Presentation,
    *,
    source_cover_index: int | None,
) -> list[PPTTemplateLayout]:
    manifests: list[PPTTemplateLayout] = []
    for index, slide in enumerate(presentation.slides):
        text_shapes = [
            shape
            for shape in _iter_source_text_shapes(slide.shapes)
            if shape.text.strip()
        ]
        content_text_shapes = [
            shape
            for shape in text_shapes
            if not any(
                term in " ".join(shape.text.split()).casefold()
                for term in ("back to contents", "목차로", "처음으로")
            )
        ]
        normalized_text = " ".join(shape.text.strip() for shape in text_shapes).casefold()
        if not normalized_text or not content_text_shapes:
            continue
        if any(term in normalized_text for term in _SOURCE_AUTHORING_EXCLUSION_TERMS):
            continue
        if _source_slide_has_chart(slide):
            continue

        has_table = _source_slide_has_table(slide)
        title = _source_slide_title(
            presentation,
            content_text_shapes,
        )
        content_capacity = _source_slide_content_capacity(
            presentation,
            content_text_shapes,
            has_table=has_table,
        )
        roles = _source_slide_roles(
            index=index,
            text=normalized_text,
            title=title.casefold(),
            source_cover_index=source_cover_index,
            has_table=has_table,
        )
        if not roles:
            continue
        role_label = _SOURCE_ROLE_LABELS[roles[0]]
        placeholder_types = ["SOURCE_TITLE"]
        if len(content_text_shapes) > 1:
            placeholder_types.append("SOURCE_BODY")
        if has_table:
            placeholder_types.append("SOURCE_TABLE")
        manifests.append(
            PPTTemplateLayout(
                layout_index=SOURCE_SLIDE_LAYOUT_OFFSET + index,
                name=f"{role_label} · {title}",
                placeholder_count=len(content_text_shapes) + int(has_table),
                placeholder_types=placeholder_types,
                supports_title=bool(content_text_shapes),
                supports_body=content_capacity > 0,
                body_placeholder_count=max(0, len(content_text_shapes) - 1)
                + int(has_table),
                source_slide_index=index,
                suggested_roles=roles,
                content_capacity=content_capacity,
            )
        )
    return manifests


def _source_slide_roles(
    *,
    index: int,
    text: str,
    title: str,
    source_cover_index: int | None,
    has_table: bool,
) -> list[str]:
    if index == source_cover_index:
        return ["cover"]
    if any(
        term in text
        for term in (
            "references",
            "sources",
            "bibliography",
            "evidence sources",
            "근거 출처",
            "참고 문헌",
        )
    ):
        return ["sources"]
    if (
        title.strip()
        in {
            "contents",
            "summary",
            "목차",
            "목차 페이지",
            "핵심정리",
            "핵심정리 페이지",
            "introduction",
            "소개",
        }
        or any(term in text for term in ("목차 페이지", "핵심정리 페이지", "introduction"))
    ):
        return ["objectives", "ncs_coverage", "lesson"]
    if (
        any(
            term in text
            for term in (
                "quick check",
                "quiz",
                "assessment",
                "before & after",
                "비교 혹은 대조",
            )
        )
        or any(term in title for term in ("평가", "문항"))
    ):
        return ["assessment", "practice"]
    if any(
        term in text
        for term in (
            "핵심키워드",
            "skill",
            "능력",
            "평균 학점",
            "수상",
            "metric",
        )
    ):
        return ["ncs_coverage", "objectives", "assessment", "practice"]
    if any(
        term in text
        for term in (
            "timeline",
            "타임라인",
            "step1",
            "project",
            "프로젝트",
            "practice",
            "exercise",
            "icebreaker",
            "실습",
        )
    ) or any(term in text for term in ("step1", "타임라인 페이지")):
        return ["practice", "assessment", "lesson"]
    if has_table:
        return ["ncs_coverage", "assessment", "lesson"]
    return ["lesson", "practice"]


def _source_slide_title(
    presentation: Presentation,
    text_shapes: list[Any],
) -> str:
    title_shape = _source_slide_title_shape(presentation, text_shapes)
    return " ".join(title_shape.text.split())[:70] or "콘텐츠 디자인"


def _source_slide_title_shape(
    presentation: Presentation,
    text_shapes: list[Any],
):
    candidates = [
        shape
        for shape in text_shapes
        if (shape.top or 0) <= presentation.slide_height * 0.35
        and (shape.width or 0) >= presentation.slide_width * 0.1
    ]
    pool = candidates or text_shapes
    return max(
        pool,
        key=lambda shape: (
            _source_shape_max_font_size(shape),
            -(shape.top or 0),
            (shape.width or 0) * (shape.height or 0),
        ),
    )


def _source_slide_content_capacity(
    presentation: Presentation,
    text_shapes: list[Any],
    *,
    has_table: bool,
) -> int:
    if not text_shapes:
        return 0
    title_shape = _source_slide_title_shape(presentation, text_shapes)
    body_shapes = [shape for shape in text_shapes if shape is not title_shape]
    slide_area = presentation.slide_width * presentation.slide_height
    body_area_ratio = (
        sum((shape.width or 0) * (shape.height or 0) for shape in body_shapes)
        / slide_area
    )
    capacity = len(body_shapes) + round(body_area_ratio * 24)
    if has_table:
        capacity = max(capacity, 4)
    return min(8, max(0, capacity))


def _source_shape_max_font_size(shape: Any) -> float:
    sizes = [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    return max(sizes, default=0.0)


def _iter_source_text_shapes(shapes: Any):
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_source_text_shapes(nested_shapes)


def _source_slide_has_table(slide: Any) -> bool:
    return any(getattr(shape, "has_table", False) for shape in slide.shapes)


def _source_slide_has_chart(slide: Any) -> bool:
    return any(getattr(shape, "has_chart", False) for shape in slide.shapes)


def _automatic_layout_mapping(layouts: list[PPTTemplateLayout]) -> dict[str, int]:
    mapping: dict[str, int] = {}
    used_source_slides: set[int] = set()
    for semantic_type in PPT_TEMPLATE_SEMANTIC_TYPES:
        ranked = sorted(
            layouts,
            key=lambda layout: _layout_score(layout, semantic_type),
            reverse=True,
        )
        if semantic_type == "sources":
            selected = ranked[0]
        else:
            selected = next(
                (
                    layout
                    for layout in ranked
                    if layout.source_slide_index is None
                    or layout.source_slide_index not in used_source_slides
                ),
                ranked[0],
            )
        mapping[semantic_type] = selected.layout_index
        if selected.source_slide_index is not None:
            used_source_slides.add(selected.source_slide_index)
    return mapping


def _layout_score(layout: PPTTemplateLayout, semantic_type: str) -> tuple[int, int, int]:
    name = layout.name.casefold()
    score = 0
    if layout.supports_title:
        score += 4
    if layout.supports_body:
        score += 4

    if layout.source_slide_index is not None:
        score += 12
        if semantic_type != "cover":
            score += layout.content_capacity * 2
            if layout.content_capacity == 0:
                score -= 80
        if layout.suggested_roles:
            if layout.suggested_roles[0] == semantic_type:
                score += 60
            elif semantic_type in layout.suggested_roles:
                role_position = layout.suggested_roles.index(semantic_type)
                score += 40 if role_position == 1 else 25 if role_position == 2 else 15

    if semantic_type == "cover":
        if any(term in name for term in ("title", "cover", "표지", "제목")):
            score += 8
        if "SUBTITLE" in layout.placeholder_types:
            score += 5
    else:
        # 표지를 제외한 모든 슬라이드는 제목+본문 콘텐츠 레이아웃을 선호한다.
        if any(term in name for term in ("content", "body", "본문", "내용")):
            score += 5
        if semantic_type == "sources" and any(
            term in name for term in ("reference", "source", "출처", "참고")
        ):
            score += 7

    # 본문 placeholder가 정확히 1개인 레이아웃을 선호한다. 생성기가 본문 하나만 채우므로,
    # 본문이 여러 개인 레이아웃은 빈 placeholder("텍스트를 입력하십시오")가 남을 수 있다.
    body_preference = 1 if layout.body_placeholder_count == 1 else 0

    return score, body_preference, -layout.layout_index


def _sanitize_filename(value: str) -> str:
    filename = value.replace("\\", "/").rsplit("/", 1)[-1].strip()
    filename = _SAFE_FILENAME.sub("_", filename).strip(" ._")
    return filename[:120] or "lessonpack-template.pptx"


def _metadata_to_row(
    metadata: PPTTemplateMetadata,
    *,
    storage_path: str,
) -> dict[str, Any]:
    return {
        "template_id": metadata.template_id,
        "project_id": metadata.project_id,
        "storage_path": storage_path,
        "original_filename": metadata.original_filename,
        "content_hash": metadata.content_hash,
        "file_size_bytes": metadata.file_size_bytes,
        "source_slide_count": metadata.source_slide_count,
        "slide_width": metadata.slide_width,
        "slide_height": metadata.slide_height,
        "layout_manifest": [item.model_dump(mode="json") for item in metadata.layouts],
        "layout_mapping": metadata.layout_mapping,
        "warnings": metadata.warnings,
        "status": metadata.status,
        "created_at": metadata.created_at.isoformat(),
        "updated_at": metadata.updated_at.isoformat(),
    }


def _metadata_from_row(row: dict[str, Any]) -> PPTTemplateMetadata:
    return PPTTemplateMetadata.model_validate(
        {
            "template_id": row["template_id"],
            "project_id": row["project_id"],
            "original_filename": row["original_filename"],
            "content_hash": row["content_hash"],
            "file_size_bytes": row["file_size_bytes"],
            "source_slide_count": row["source_slide_count"],
            "slide_width": row["slide_width"],
            "slide_height": row["slide_height"],
            "layouts": row.get("layout_manifest", []),
            "layout_mapping": row.get("layout_mapping", {}),
            "warnings": row.get("warnings", []),
            "status": row.get("status", "ready"),
            "created_at": row["created_at"],
            "updated_at": row["updated_at"],
        }
    )


def _response_data(response: Any) -> list[dict[str, Any]]:
    error = getattr(response, "error", None)
    if error is None and isinstance(response, dict):
        error = response.get("error")
    if error:
        raise RuntimeError(f"Supabase PPT template request failed: {error}")
    data = getattr(response, "data", None)
    if data is None and isinstance(response, dict):
        data = response.get("data")
    return data if isinstance(data, list) else []


def _raise_for_supabase_error(response: Any) -> None:
    error = getattr(response, "error", None)
    if error is None and isinstance(response, dict):
        error = response.get("error")
    if error:
        raise RuntimeError(f"Supabase PPT template request failed: {error}")


def _best_effort_remove(bucket: Any, storage_path: str) -> None:
    try:
        bucket.remove([storage_path])
    except Exception:
        return
