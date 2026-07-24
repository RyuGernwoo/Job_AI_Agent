from __future__ import annotations

import io
import re
from collections import Counter, OrderedDict
from copy import deepcopy
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from lectureops_agent.models.schemas import CourseType, CitationDetail, LessonPackage, NCSAlignment
from lectureops_agent.services.ppt_template_service import (
    resolve_template_mapping_for_export,
    reusable_source_cover_index,
    source_slide_index_for_layout,
)


_INVALID_FILENAME_CHARS = re.compile(r'[<>:"/\\|?*\x00-\x1f]')
_FILENAME_SEPARATOR = re.compile(r"[\s_]+")
_SOURCE_NAVIGATION_TERMS = ("back to contents", "목차로", "처음으로")


def export_lesson_package_docx(*, package: LessonPackage, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    document.add_heading(package.lesson_plan.title, level=0)
    document.add_paragraph(f"템플릿 버전: {package.template_metadata.template_version}")
    if package.template_metadata.lesson_duration_min:
        document.add_paragraph(f"수업 시간: {package.template_metadata.lesson_duration_min}분")
    training_plan = _training_plan_summary(package)
    if training_plan:
        document.add_paragraph(training_plan)

    document.add_heading("학습 목표", level=1)
    for objective in package.lesson_plan.learning_objectives:
        document.add_paragraph(objective, style="List Bullet")
    if package.course_type == CourseType.GENERAL:
        document.add_heading("학습목표-활동-평가 연결", level=1)
        for objective in package.lesson_plan.learning_objectives:
            document.add_paragraph(
                f"{objective} → 교안 설명 → 실습 수행 → 객관식·수행평가 확인",
                style="List Bullet",
            )

    document.add_heading("교안", level=1)
    for flow in package.lesson_plan.lecture_flow:
        document.add_heading(flow.section, level=2)
        if flow.duration_min:
            document.add_paragraph(f"예상 시간: {flow.duration_min}분")
        document.add_paragraph(flow.content)
        _add_alignment_paragraph(document, flow.ncs_alignment)

    document.add_heading("실습 과제", level=1)
    document.add_paragraph(_without_label(package.practice.scenario, "실습 시나리오"))
    document.add_heading("수행 절차", level=2)
    for step in package.practice.steps:
        document.add_paragraph(_without_label(step, "수행 절차"), style="List Number")
    document.add_heading("제출물", level=2)
    document.add_paragraph(_without_label(package.practice.submission, "제출물"))
    document.add_heading("실습 루브릭", level=2)
    for item in package.practice.rubric:
        document.add_paragraph(_without_label(item, "평가 기준"), style="List Bullet")
    _add_alignment_paragraph(document, package.practice.ncs_alignment)

    document.add_heading("평가 문항", level=1)
    for index, question in enumerate(package.assessment.multiple_choice, start=1):
        document.add_paragraph(f"문항 {index}. {question.question}")
        for option_index, option in enumerate(question.options, start=1):
            document.add_paragraph(f"{option_index}. {option}", style="List Bullet")
        document.add_paragraph(f"정답: {question.answer_index + 1}")
        document.add_paragraph(f"해설: {question.explanation}")
        _add_alignment_paragraph(document, question.ncs_alignment)

    task = package.assessment.performance_task
    document.add_heading("수행평가", level=2)
    document.add_paragraph(task.title)
    document.add_paragraph(task.description)
    for item in task.rubric:
        document.add_paragraph(_without_label(item, "평가 기준"), style="List Bullet")
    _add_alignment_paragraph(document, task.ncs_alignment)

    _add_ncs_coverage_section(document, package)
    _add_evidence_section(document, package)

    document.save(output_path)
    return output_path


def export_lesson_package_pptx(
    *,
    package: LessonPackage,
    output_path: Path,
    template_content: bytes | None = None,
    layout_mapping: dict[str, int] | None = None,
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = (
        Presentation(io.BytesIO(template_content))
        if template_content is not None
        else Presentation()
    )
    source_slide_count = len(presentation.slides) if template_content is not None else 0
    mapping = (
        resolve_template_mapping_for_export(presentation, layout_mapping)
        if template_content is not None
        else dict(layout_mapping or {})
    )
    preserved_brand_texts = _repeated_brand_texts(
        presentation,
        source_slide_count=source_slide_count,
    )
    source_cover_index = source_slide_index_for_layout(mapping.get("cover"))
    if source_cover_index is None and template_content is not None:
        source_cover_index = reusable_source_cover_index(presentation)

    # 표지
    training_plan = _training_plan_summary(package)
    if (
        source_cover_index is not None
        and 0 <= source_cover_index < source_slide_count
    ):
        cover = _duplicate_source_slide(
            presentation,
            presentation.slides[source_cover_index],
        )
        _populate_reused_source_slide(
            presentation,
            cover,
            title=package.lesson_plan.title,
            bullets=[
                (value, 0)
                for value in ["LessonPack AI 강의 패키지", training_plan]
                if value
            ],
            preserved_texts=preserved_brand_texts,
            is_cover=True,
        )
    else:
        cover = presentation.slides.add_slide(
            _resolve_layout(presentation, mapping, "cover", fallback_index=0)
        )
        _set_slide_title(presentation, cover, package.lesson_plan.title, is_cover=True)
        _set_slide_bullets(
            presentation,
            cover,
            [(value, 0) for value in ["LessonPack AI 강의 패키지", training_plan] if value],
            prefer_subtitle=True,
        )

    # 학습 목표
    _add_bullet_slides(
        presentation,
        "학습 목표",
        [(objective, 0) for objective in package.lesson_plan.learning_objectives],
        semantic_type="objectives",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )
    if package.course_type == CourseType.GENERAL:
        _add_bullet_slides(
            presentation,
            "학습목표 · 활동 · 평가 연결",
            [
                (f"{objective} → 교안 → 실습 → 평가", 0)
                for objective in package.lesson_plan.learning_objectives
            ],
            semantic_type="objectives",
            layout_mapping=mapping,
            source_slide_count=source_slide_count,
            preserved_texts=preserved_brand_texts,
        )

    # 교안 (섹션별, 본문을 문장 단위 불릿으로 분해해 가독성 확보)
    for flow in package.lesson_plan.lecture_flow:
        blocks: list[tuple[str, int]] = []
        if flow.duration_min:
            blocks.append((f"예상 시간 {flow.duration_min}분", 0))
        blocks.extend((sentence, 0) for sentence in _sentence_bullets(flow.content))
        blocks.extend((item, 0) for item in _alignment_bullets(flow.ncs_alignment))
        _add_bullet_slides(
            presentation,
            f"교안 · {flow.section}",
            blocks,
            semantic_type="lesson",
            layout_mapping=mapping,
            source_slide_count=source_slide_count,
            preserved_texts=preserved_brand_texts,
        )

    # 실습
    practice = package.practice
    _add_bullet_slides(
        presentation,
        "실습 개요",
        [
            *[(s, 0) for s in _sentence_bullets(_without_label(practice.scenario, "실습 시나리오"))],
            (f"제출물: {_without_label(practice.submission, '제출물')}", 0),
        ],
        semantic_type="practice",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )
    _add_bullet_slides(
        presentation,
        "실습 수행 절차",
        [
            (f"{index}. {_without_label(step, '수행 절차')}", 0)
            for index, step in enumerate(practice.steps, start=1)
        ],
        semantic_type="practice",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )
    _add_bullet_slides(
        presentation,
        "실습 평가 기준(루브릭)",
        [
            *[(_without_label(item, "평가 기준"), 0) for item in practice.rubric],
            *[(item, 0) for item in _alignment_bullets(practice.ncs_alignment)],
        ],
        semantic_type="practice",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )

    # 평가 개요
    task = package.assessment.performance_task
    _add_bullet_slides(
        presentation,
        "평가 개요",
        [
            (f"객관식 문항 {len(package.assessment.multiple_choice)}개 · 수행평가 1개", 0),
            (f"수행평가: {task.title}", 0),
            *[(s, 0) for s in _sentence_bullets(task.description)],
        ],
        semantic_type="assessment",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )
    # 객관식 전 문항: 보기(선택지)와 정답·해설을 한 문항당 한 슬라이드로 제시한다.
    for index, question in enumerate(package.assessment.multiple_choice, start=1):
        blocks = [(question.question, 0)]
        blocks.extend(
            (f"{option_index}. {option}", 1)
            for option_index, option in enumerate(question.options, start=1)
        )
        blocks.append((f"정답: {question.answer_index + 1}번", 0))
        blocks.append((f"해설: {question.explanation}", 0))
        _add_bullet_slides(
            presentation,
            f"평가 문항 {index}",
            blocks,
            semantic_type="assessment",
            layout_mapping=mapping,
            paginate=False,
            source_slide_count=source_slide_count,
            preserved_texts=preserved_brand_texts,
        )

    # 수행평가 상세
    _add_bullet_slides(
        presentation,
        "수행평가",
        [
            (task.title, 0),
            *[(s, 0) for s in _sentence_bullets(task.description)],
            *[(_without_label(item, "평가 기준"), 0) for item in task.rubric],
            *[(item, 0) for item in _alignment_bullets(task.ncs_alignment)],
        ],
        semantic_type="assessment",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )

    if package.ncs_coverage is not None:
        _add_bullet_slides(
            presentation,
            "NCS 수행준거 커버리지",
            [(bullet, 0) for bullet in _ncs_coverage_bullets(package)],
            semantic_type="ncs_coverage",
            layout_mapping=mapping,
            source_slide_count=source_slide_count,
            preserved_texts=preserved_brand_texts,
        )

    # 근거는 발표 흐름을 방해하지 않도록 마지막 출처 슬라이드에만 표시한다.
    _add_bullet_slides(
        presentation,
        "근거 출처",
        [(bullet, 0) for bullet in _compact_evidence_bullets(package)],
        semantic_type="sources",
        layout_mapping=mapping,
        source_slide_count=source_slide_count,
        preserved_texts=preserved_brand_texts,
    )

    if source_slide_count:
        _remove_source_slides(presentation, source_slide_count=source_slide_count)
    presentation.save(output_path)
    _validate_exported_pptx(output_path)
    return output_path


def build_export_filename(package: LessonPackage, extension: str) -> str:
    normalized_extension = extension.casefold().lstrip(".")
    if normalized_extension not in {"docx", "pptx"}:
        raise ValueError(f"unsupported export extension: {extension}")

    title = _INVALID_FILENAME_CHARS.sub(" ", package.lesson_plan.title)
    title = _FILENAME_SEPARATOR.sub("_", title).strip("._ ")
    title = title[:48].rstrip("._ ") or "LessonPack"
    suffix = "교안" if normalized_extension == "docx" else "강의자료"
    return f"{title}_{suffix}.{normalized_extension}"


def _training_plan_summary(package: LessonPackage) -> str:
    metadata = package.template_metadata
    if (
        metadata.total_training_hours is None
        or metadata.total_lessons is None
        or metadata.theory_ratio_percent is None
        or metadata.practice_ratio_percent is None
    ):
        return ""
    return (
        f"훈련 운영: 총 {metadata.total_training_hours:g}시간 · {metadata.total_lessons}차시 · "
        f"이론 {metadata.theory_ratio_percent}% · 실습 {metadata.practice_ratio_percent}%"
    )


def _add_alignment_paragraph(document: Document, alignments: list[NCSAlignment]) -> None:
    if not alignments:
        return
    for alignment in alignments:
        document.add_paragraph(f"NCS 연계: {alignment.unit_code} {alignment.unit_name}")


def _add_ncs_coverage_section(document: Document, package: LessonPackage) -> None:
    report = package.ncs_coverage
    if report is None:
        return
    document.add_heading("NCS 수행준거 커버리지", level=1)
    document.add_paragraph(
        f"설계 커버리지 {report.covered_criteria_count}/{report.target_criteria_count} "
        f"({report.coverage * 100:.0f}%), 평가 커버리지 "
        f"{report.assessment_criteria_count}/{report.target_criteria_count} "
        f"({report.assessment_coverage * 100:.0f}%)"
    )
    for item in report.items:
        locations = [*item.lesson_sections]
        if item.practice:
            locations.append("실습")
        locations.extend(item.assessment_items)
        document.add_paragraph(
            f"{item.unit_code} {item.performance_criterion} | "
            f"연결: {', '.join(locations) if locations else '없음'}",
            style="List Bullet",
        )
    for warning in report.warnings:
        document.add_paragraph(f"주의: {warning}")


def _add_evidence_section(document: Document, package: LessonPackage) -> None:
    document.add_heading("근거 출처", level=1)
    document.add_paragraph("본문에서 사용한 근거를 출처별로 정리했습니다.")
    for item in _grouped_evidence_bullets(package, include_url=True):
        document.add_paragraph(item, style="List Bullet")


def _collect_citation_ids(package: LessonPackage) -> list[str]:
    ordered: list[str] = []
    seen: set[str] = set()

    def add_many(citation_ids: list[str]) -> None:
        for citation_id in citation_ids:
            if citation_id not in seen:
                seen.add(citation_id)
                ordered.append(citation_id)

    for flow in package.lesson_plan.lecture_flow:
        add_many(flow.citation_ids)
    add_many(package.practice.citation_ids)
    for question in package.assessment.multiple_choice:
        add_many(question.citation_ids)
    add_many(package.assessment.performance_task.citation_ids)
    return ordered


def _evidence_details(package: LessonPackage) -> list[CitationDetail]:
    citation_ids = _collect_citation_ids(package)
    details_by_id = {detail.chunk_id: detail for detail in package.evidence_sources}
    return [
        details_by_id.get(
            citation_id,
            CitationDetail(
                chunk_id=citation_id,
                source_name="출처 정보 미기재",
                excerpt="원문 발췌 정보가 패키지에 저장되지 않았습니다.",
            ),
        )
        for citation_id in citation_ids
    ]


def _alignment_bullets(alignments: list[NCSAlignment]) -> list[str]:
    if not alignments:
        return []
    return [f"NCS 연계: {alignment.unit_code} {alignment.unit_name}" for alignment in alignments]


def _ncs_coverage_bullets(package: LessonPackage) -> list[str]:
    report = package.ncs_coverage
    if report is None:
        return []
    bullets = [
        f"설계 커버리지: {report.covered_criteria_count}/{report.target_criteria_count} "
        f"({report.coverage * 100:.0f}%)",
        f"평가 커버리지: {report.assessment_criteria_count}/{report.target_criteria_count} "
        f"({report.assessment_coverage * 100:.0f}%)",
    ]
    bullets.extend(
        f"{item.unit_code}: {item.performance_criterion}"
        for item in report.items[:3]
    )
    bullets.extend(f"주의: {warning}" for warning in report.warnings)
    return bullets


def _compact_evidence_bullets(package: LessonPackage) -> list[str]:
    return _grouped_evidence_bullets(package, include_url=False)


def _grouped_evidence_bullets(package: LessonPackage, *, include_url: bool) -> list[str]:
    grouped: OrderedDict[str, dict[str, object]] = OrderedDict()
    for detail in _evidence_details(package):
        group = grouped.setdefault(
            detail.source_name,
            {"chunk_ids": [], "source_url": detail.source_url, "license": detail.license},
        )
        chunk_ids = group["chunk_ids"]
        assert isinstance(chunk_ids, list)
        chunk_ids.append(detail.chunk_id)
        if not group["source_url"] and detail.source_url:
            group["source_url"] = detail.source_url
        if not group["license"] and detail.license:
            group["license"] = detail.license

    bullets: list[str] = []
    for source_name, values in grouped.items():
        chunk_ids = ", ".join(str(value) for value in values["chunk_ids"])
        text = f"{source_name} | 사용 청크: {chunk_ids}"
        if include_url and values["source_url"]:
            text += f" | URL: {values['source_url']}"
        if values["license"]:
            text += f" | 라이선스: {values['license']}"
        elif not include_url and values["source_url"]:
            text += f" | {values['source_url']}"
        bullets.append(text)
    return bullets or ["사용된 근거 출처가 없습니다."]


def _without_label(value: str, label: str) -> str:
    normalized = value.strip()
    prefix = f"{label}:"
    if normalized.startswith(prefix):
        return normalized[len(prefix) :].strip()
    return normalized


# 한 슬라이드에 담을 최대 블록 수. autofit과 함께 가독성을 유지한다.
_SLIDE_BLOCK_CAPACITY = 8
_BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "SUBTITLE", "VERTICAL_BODY", "VERTICAL_OBJECT"}
# 채우지 않으면 "텍스트를 입력하십시오" 안내가 남는 placeholder 유형 (슬라이드에서 제거한다).
_REMOVABLE_PLACEHOLDER_TYPES = _BODY_PLACEHOLDER_TYPES | {"PICTURE", "CONTENT"}
_SENTENCE_BOUNDARY = re.compile(r"(?<=[.!?])\s+|(?<=다\.)\s+|(?<=요\.)\s+|(?<=음\.)\s+")

Block = tuple[str, int]


def _sentence_bullets(text: str, *, max_len: int = 220) -> list[str]:
    """긴 본문 문단을 문장 단위 불릿으로 분해해 슬라이드 가독성을 높인다."""
    normalized = str(text or "").strip()
    if not normalized:
        return []
    bullets: list[str] = []
    for part in re.split(r"\n+", normalized):
        part = part.strip()
        if not part:
            continue
        buffer = ""
        for sentence in _SENTENCE_BOUNDARY.split(part):
            sentence = sentence.strip()
            if not sentence:
                continue
            if buffer and len(buffer) + len(sentence) + 1 > max_len:
                bullets.append(buffer)
                buffer = sentence
            else:
                buffer = f"{buffer} {sentence}".strip()
        if buffer:
            bullets.append(buffer)
    return bullets


def _normalize_blocks(blocks: list) -> list[Block]:
    normalized: list[Block] = []
    for block in blocks:
        text, level = block if isinstance(block, tuple) else (block, 0)
        text = " ".join(str(text).split())
        if text:
            normalized.append((text, int(level)))
    return normalized


def _add_bullet_slides(
    presentation: Presentation,
    title: str,
    bullets: list,
    *,
    semantic_type: str,
    layout_mapping: dict[str, int],
    paginate: bool = True,
    source_slide_count: int = 0,
    preserved_texts: set[str] | None = None,
) -> None:
    normalized = _normalize_blocks(bullets)
    if not normalized:
        normalized = [("내용이 없습니다.", 0)]
    source_slide_index = source_slide_index_for_layout(
        layout_mapping.get(semantic_type)
    )
    block_capacity = _SLIDE_BLOCK_CAPACITY
    if (
        source_slide_index is not None
        and 0 <= source_slide_index < source_slide_count
    ):
        block_capacity = _source_slide_block_capacity(
            presentation,
            presentation.slides[source_slide_index],
            preserved_texts=preserved_texts or set(),
        )
    if paginate:
        groups = [
            normalized[index : index + block_capacity]
            for index in range(0, len(normalized), block_capacity)
        ]
    else:
        groups = [normalized]
    for index, group in enumerate(groups):
        continued_title = title if index == 0 else f"{title} (계속)"
        if (
            source_slide_index is not None
            and 0 <= source_slide_index < source_slide_count
        ):
            slide = _duplicate_source_slide(
                presentation,
                presentation.slides[source_slide_index],
            )
            _populate_reused_source_slide(
                presentation,
                slide,
                title=continued_title,
                bullets=group,
                preserved_texts=preserved_texts or set(),
            )
        else:
            slide = presentation.slides.add_slide(
                _resolve_layout(
                    presentation,
                    layout_mapping,
                    semantic_type,
                    fallback_index=1,
                )
            )
            _set_slide_title(presentation, slide, continued_title)
            _set_slide_bullets(presentation, slide, group)


def _source_slide_block_capacity(
    presentation: Presentation,
    slide,
    *,
    preserved_texts: set[str],
) -> int:
    text_shapes = [
        shape
        for shape in _iter_text_shapes(slide.shapes)
        if shape.text.strip()
        and _normalize_source_text(shape.text) not in preserved_texts
        and not any(
            term in _normalize_source_text(shape.text)
            for term in _SOURCE_NAVIGATION_TERMS
        )
    ]
    if len(text_shapes) <= 1:
        return 1
    title_shape = _select_source_title_shape(presentation, text_shapes)
    body_shapes = [shape for shape in text_shapes if shape is not title_shape]
    slide_area = presentation.slide_width * presentation.slide_height
    body_area_ratio = (
        sum((shape.width or 0) * (shape.height or 0) for shape in body_shapes)
        / slide_area
    )
    capacity = len(body_shapes) + round(body_area_ratio * 24)
    if any(getattr(shape, "has_table", False) for shape in slide.shapes):
        capacity = max(capacity, 4)
    return min(_SLIDE_BLOCK_CAPACITY, max(1, capacity))


def _set_slide_title(
    presentation: Presentation,
    slide,
    title: str,
    *,
    is_cover: bool = False,
) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = _truncate(title, max_length=100)
        return
    box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(0.55),
        presentation.slide_width - Inches(1.4),
        Inches(1.0),
    )
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = _truncate(title, max_length=100)
    paragraph.font.size = Pt(30 if is_cover else 26)
    paragraph.font.bold = True


def _set_slide_bullets(
    presentation: Presentation,
    slide,
    bullets: list,
    *,
    prefer_subtitle: bool = False,
) -> None:
    body = _prepare_body_frame(slide, prefer_subtitle=prefer_subtitle)
    if body is None:
        box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.6),
            presentation.slide_width - Inches(1.6),
            presentation.slide_height - Inches(2.2),
        )
        body = box.text_frame
    body.word_wrap = True
    try:
        # 본문이 넘칠 경우 자동으로 글자 크기를 줄여 잘림을 방지한다.
        body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except (ValueError, NotImplementedError):
        pass
    body.clear()
    for index, (text, level) in enumerate(_normalize_blocks(bullets)):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = _truncate(text, max_length=500)
        paragraph.level = min(max(level, 0), 4)
        # 템플릿의 과도하게 큰 본문 기본값이 content box를 넘지 않도록 상한을 둔다.
        paragraph.font.size = Pt(18 if level == 0 else 16)


def _prepare_body_frame(slide, *, prefer_subtitle: bool):
    """본문 placeholder를 하나 선택하고, 나머지 빈 본문 placeholder는 제거해 깔끔한 슬라이드를 만든다."""
    title_shape = slide.shapes.title
    candidates: list[tuple[str, object]] = []
    for placeholder in list(slide.placeholders):
        if title_shape is not None and placeholder is title_shape:
            continue
        placeholder_type = getattr(
            placeholder.placeholder_format.type,
            "name",
            str(placeholder.placeholder_format.type),
        )
        candidates.append((placeholder_type, placeholder))

    preferred = (
        {"SUBTITLE"}
        if prefer_subtitle
        else {"BODY", "OBJECT", "VERTICAL_BODY", "VERTICAL_OBJECT"}
    )
    chosen = None
    for placeholder_type, placeholder in candidates:
        if placeholder.has_text_frame and placeholder_type in preferred:
            chosen = placeholder
            break
    if chosen is None:
        for placeholder_type, placeholder in candidates:
            if placeholder.has_text_frame and placeholder_type in _BODY_PLACEHOLDER_TYPES:
                chosen = placeholder
                break

    for placeholder_type, placeholder in candidates:
        if placeholder is chosen:
            continue
        if placeholder_type in _REMOVABLE_PLACEHOLDER_TYPES:
            element = placeholder._element
            element.getparent().remove(element)

    return chosen.text_frame if chosen is not None else None


def _resolve_layout(
    presentation: Presentation,
    layout_mapping: dict[str, int],
    semantic_type: str,
    *,
    fallback_index: int,
):
    if not presentation.slide_layouts:
        raise ValueError("PPT presentation does not contain any slide layouts")
    layout_index = layout_mapping.get(semantic_type, fallback_index)
    if layout_index < 0 or layout_index >= len(presentation.slide_layouts):
        layout_index = min(fallback_index, len(presentation.slide_layouts) - 1)
    return presentation.slide_layouts[layout_index]


def _remove_source_slides(
    presentation: Presentation,
    *,
    source_slide_count: int,
) -> None:
    slide_ids = presentation.slides._sldIdLst
    source_ids = list(slide_ids)[:source_slide_count]
    for slide_id in reversed(source_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def _duplicate_source_slide(
    presentation: Presentation,
    source_slide,
):
    destination = presentation.slides.add_slide(source_slide.slide_layout)
    relationship_map: dict[str, str] = {}
    skipped_relationships: set[str] = set()
    for relationship in source_slide.part.rels.values():
        if (
            relationship.is_external
            or relationship.reltype.endswith(("/slideLayout", "/notesSlide", "/slide"))
        ):
            skipped_relationships.add(relationship.rId)
            continue
        relationship_map[relationship.rId] = destination.part.relate_to(
            relationship.target_part,
            relationship.reltype,
        )

    destination_tree = destination.shapes._spTree
    for element in list(destination_tree)[2:]:
        destination_tree.remove(element)
    for source_element in list(source_slide.shapes._spTree)[2:]:
        if source_element.tag.endswith("}extLst"):
            continue
        copied_element = deepcopy(source_element)
        _remap_copied_relationships(
            copied_element,
            relationship_map=relationship_map,
            skipped_relationships=skipped_relationships,
        )
        destination_tree.insert_element_before(copied_element, "p:extLst")

    destination_c_sld = destination._element.cSld
    for element in list(destination_c_sld):
        if element.tag.endswith("}bg"):
            destination_c_sld.remove(element)
    source_background = next(
        (
            element
            for element in source_slide._element.cSld
            if element.tag.endswith("}bg")
        ),
        None,
    )
    if source_background is not None:
        copied_background = deepcopy(source_background)
        _remap_copied_relationships(
            copied_background,
            relationship_map=relationship_map,
            skipped_relationships=skipped_relationships,
        )
        destination_c_sld.insert(0, copied_background)
    return destination


def _remap_copied_relationships(
    element,
    *,
    relationship_map: dict[str, str],
    skipped_relationships: set[str],
) -> None:
    removable_nodes = []
    for node in element.iter():
        for attribute_name, attribute_value in list(node.attrib.items()):
            if attribute_value in relationship_map:
                node.set(attribute_name, relationship_map[attribute_value])
            elif attribute_value in skipped_relationships:
                if node.tag.endswith(("}hlinkClick", "}hlinkHover")):
                    removable_nodes.append(node)
                else:
                    node.attrib.pop(attribute_name, None)
    for node in removable_nodes:
        if node.getparent() is not None:
            node.getparent().remove(node)


def _repeated_brand_texts(
    presentation: Presentation,
    *,
    source_slide_count: int,
) -> set[str]:
    counts: Counter[str] = Counter()
    for slide in list(presentation.slides)[:source_slide_count]:
        slide_values = {
            _normalize_source_text(shape.text)
            for shape in _iter_text_shapes(slide.shapes)
            if shape.text.strip()
        }
        counts.update(slide_values)
    repeated_threshold = max(3, (source_slide_count * 3 + 4) // 5)
    return {
        value
        for value, count in counts.items()
        if count >= repeated_threshold
        and len(value) <= 80
        and not any(term in value for term in _SOURCE_NAVIGATION_TERMS)
    }


def _normalize_source_text(value: str) -> str:
    return " ".join(value.split()).casefold()


def _populate_reused_source_slide(
    presentation: Presentation,
    slide,
    *,
    title: str,
    bullets: list[Block],
    preserved_texts: set[str],
    is_cover: bool = False,
) -> None:
    text_shapes = [
        shape
        for shape in _iter_text_shapes(slide.shapes)
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]
    editable_shapes = [
        shape
        for shape in text_shapes
        if _normalize_source_text(shape.text) not in preserved_texts
    ]
    if not editable_shapes:
        _set_slide_title(presentation, slide, title, is_cover=True)
        _set_slide_bullets(
            presentation,
            slide,
            bullets,
            prefer_subtitle=is_cover,
        )
        return

    title_shape = _select_source_title_shape(presentation, editable_shapes)
    body_shapes = sorted(
        [shape for shape in editable_shapes if shape is not title_shape],
        key=lambda shape: (
            (shape.width or 0) * (shape.height or 0),
            _maximum_text_font_size(shape),
        ),
        reverse=True,
    )
    slide_area = presentation.slide_width * presentation.slide_height

    for shape in editable_shapes:
        _replace_shape_text_preserving_style(shape, "")
    for shape in text_shapes:
        if any(
            term in _normalize_source_text(shape.text)
            for term in _SOURCE_NAVIGATION_TERMS
        ):
            _replace_shape_text_preserving_style(shape, "")

    _ensure_source_title_bounds(presentation, title_shape)
    _replace_shape_text_preserving_style(
        title_shape,
        _truncate(title, max_length=100),
        max_font_size=_source_title_font_cap(
            presentation,
            title_shape,
            title,
            is_cover=is_cover,
        ),
    )

    normalized = _normalize_blocks(bullets)
    if is_cover:
        _separate_source_title_and_body(
            presentation,
            title_shape,
            body_shapes,
        )
        subtitle_shape = body_shapes[0] if body_shapes else None
        subtitle = "\n".join(text for text, _ in normalized)
        if subtitle_shape is not None:
            _replace_shape_text_preserving_style(
                subtitle_shape,
                _truncate(subtitle, max_length=180),
                max_font_size=16,
            )
        elif subtitle:
            _set_slide_bullets(
                presentation,
                slide,
                [(subtitle, 0)],
                prefer_subtitle=True,
            )
        _clear_source_tables(slide)
        return

    tables = [
        shape.table
        for shape in slide.shapes
        if getattr(shape, "has_table", False)
    ]
    if tables:
        _populate_source_table(tables[0], normalized)
        for table in tables[1:]:
            _populate_source_table(table, [])
        return

    display_lines = [
        f"{'  ' * max(level, 0)}• {_truncate(text, max_length=220)}"
        for text, level in normalized
    ]
    usable_body_shapes = [
        shape
        for shape in body_shapes
        if ((shape.width or 0) * (shape.height or 0)) / slide_area >= 0.01
    ] or body_shapes[:1]
    if (
        len(display_lines) > max(2, len(usable_body_shapes) * 2)
        and usable_body_shapes
        and _can_merge_source_body_zone(presentation, body_shapes)
    ):
        _expand_shape_to_union(usable_body_shapes[0], body_shapes)
        usable_body_shapes = usable_body_shapes[:1]

    _separate_source_title_and_body(
        presentation,
        title_shape,
        usable_body_shapes,
    )
    slot_count = min(len(usable_body_shapes), len(display_lines), 4)
    if slot_count:
        for slot_index, shape in enumerate(usable_body_shapes[:slot_count]):
            start = (len(display_lines) * slot_index) // slot_count
            end = (len(display_lines) * (slot_index + 1)) // slot_count
            value = _fit_source_body_copy(
                presentation,
                shape,
                "\n".join(display_lines[start:end]),
            )
            _replace_shape_text_preserving_style(
                shape,
                value,
                max_font_size=_source_body_font_cap(
                    presentation,
                    shape,
                    value,
                ),
            )
    elif normalized:
        _set_slide_bullets(
            presentation,
            slide,
            normalized,
        )


def _source_title_font_cap(
    presentation: Presentation,
    shape,
    title: str,
    *,
    is_cover: bool,
) -> int:
    width_ratio = (shape.width or 0) / presentation.slide_width
    cap = 42 if is_cover else 34
    if width_ratio < 0.25:
        cap = min(cap, 24)
    elif width_ratio < 0.4 or len(title) > 18:
        cap = min(cap, 28)
    if len(title) > 36:
        cap = min(cap, 24)
    return cap


def _ensure_source_title_bounds(
    presentation: Presentation,
    shape,
) -> None:
    minimum_width = int(presentation.slide_width * 0.32)
    minimum_height = int(presentation.slide_height * 0.09)
    if (shape.width or 0) < minimum_width:
        center = (shape.left or 0) + (shape.width or 0) // 2
        shape.width = minimum_width
        shape.left = max(
            0,
            min(
                center - minimum_width // 2,
                presentation.slide_width - minimum_width,
            ),
        )
    if (shape.height or 0) < minimum_height:
        shape.height = minimum_height


def _separate_source_title_and_body(
    presentation: Presentation,
    title_shape,
    body_shapes: list,
) -> None:
    title_left = title_shape.left or 0
    title_right = title_left + (title_shape.width or 0)
    title_bottom = (title_shape.top or 0) + (title_shape.height or 0)
    margin = int(presentation.slide_height * 0.015)
    maximum_bottom = int(presentation.slide_height * 0.9)

    for shape in body_shapes:
        shape_left = shape.left or 0
        shape_right = shape_left + (shape.width or 0)
        shape_top = shape.top or 0
        shape_bottom = shape_top + (shape.height or 0)
        horizontally_overlaps = min(title_right, shape_right) > max(
            title_left,
            shape_left,
        )
        vertically_overlaps = min(title_bottom, shape_bottom) > max(
            title_shape.top or 0,
            shape_top,
        )
        if not horizontally_overlaps or not vertically_overlaps:
            continue

        new_top = title_bottom + margin
        available_height = maximum_bottom - new_top
        if available_height <= 0:
            continue
        shape.top = new_top
        shape.height = min(shape.height or available_height, available_height)


def _select_source_title_shape(
    presentation: Presentation,
    text_shapes: list,
):
    candidates = [
        shape
        for shape in text_shapes
        if (shape.top or 0) <= presentation.slide_height * 0.35
        and (shape.width or 0) >= presentation.slide_width * 0.1
    ]
    pool = candidates or text_shapes
    return max(
        pool,
        key=lambda shape: (
            _maximum_text_font_size(shape),
            -(shape.top or 0),
            (shape.width or 0) * (shape.height or 0),
        ),
    )


def _source_body_font_cap(
    presentation: Presentation,
    shape,
    value: str,
) -> int:
    area_ratio = (
        (shape.width or 0) * (shape.height or 0)
        / (presentation.slide_width * presentation.slide_height)
    )
    height_ratio = (shape.height or 0) / presentation.slide_height
    line_count = value.count("\n") + 1
    if area_ratio < 0.015:
        return 11
    if area_ratio < 0.03:
        return 13
    if height_ratio < 0.15 and len(value) > 70:
        return 11
    if line_count >= 5 or len(value) > 320:
        return 13
    if (
        line_count >= 3
        or (line_count >= 2 and area_ratio < 0.06)
        or area_ratio < 0.04
        or len(value) > 180
    ):
        return 15
    return 18


def _fit_source_body_copy(
    presentation: Presentation,
    shape,
    value: str,
) -> str:
    area_ratio = (
        (shape.width or 0) * (shape.height or 0)
        / (presentation.slide_width * presentation.slide_height)
    )
    if area_ratio < 0.015:
        return _truncate(value, max_length=65)
    if area_ratio < 0.03:
        return _truncate(value, max_length=120)
    return value


def _can_merge_source_body_zone(
    presentation: Presentation,
    body_shapes: list,
) -> bool:
    if not body_shapes:
        return False
    left = min(shape.left or 0 for shape in body_shapes)
    right = max((shape.left or 0) + (shape.width or 0) for shape in body_shapes)
    width_ratio = (right - left) / presentation.slide_width
    return width_ratio <= 0.55


def _expand_shape_to_union(shape, shapes: list) -> None:
    left = min(item.left or 0 for item in shapes)
    top = min(item.top or 0 for item in shapes)
    right = max((item.left or 0) + (item.width or 0) for item in shapes)
    bottom = max((item.top or 0) + (item.height or 0) for item in shapes)
    shape.left = left
    shape.top = top
    shape.width = right - left
    shape.height = bottom - top


def _clear_source_tables(slide) -> None:
    for shape in list(slide.shapes):
        if getattr(shape, "has_table", False):
            element = shape._element
            element.getparent().remove(element)


def _populate_source_table(table, blocks: list[Block]) -> None:
    values = [_truncate(text, max_length=160) for text, _ in blocks]
    cells = [cell for row in table.rows for cell in row.cells]
    for index, cell in enumerate(cells):
        value = values[index] if index < len(values) else ""
        _replace_text_frame_preserving_style(
            cell.text_frame,
            value,
            max_font_size=14,
        )


def _maximum_text_font_size(shape) -> float:
    sizes = [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    return max(sizes, default=0.0)


def _iter_text_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_text_shapes(nested_shapes)


def _replace_shape_text_preserving_style(
    shape,
    value: str,
    *,
    max_font_size: int | None = None,
) -> None:
    _replace_text_frame_preserving_style(
        shape.text_frame,
        value,
        max_font_size=max_font_size,
    )


def _replace_text_frame_preserving_style(
    text_frame,
    value: str,
    *,
    max_font_size: int | None = None,
) -> None:
    text_frame.word_wrap = True
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except (ValueError, NotImplementedError):
        pass
    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.line_spacing = 1.0
    first_paragraph.space_before = Pt(0)
    first_paragraph.space_after = Pt(0)
    first_runs = list(first_paragraph.runs)
    for run in first_paragraph.runs:
        run.text = ""
    for paragraph in list(text_frame.paragraphs)[1:]:
        paragraph._p.getparent().remove(paragraph._p)
    if first_runs:
        first_runs[0].text = value
        if max_font_size is not None and (
            first_runs[0].font.size is None
            or first_runs[0].font.size.pt > max_font_size
        ):
            first_runs[0].font.size = Pt(max_font_size)
    else:
        first_paragraph.text = value
        if max_font_size is not None:
            first_paragraph.font.size = Pt(max_font_size)


def _validate_exported_pptx(output_path: Path) -> None:
    try:
        exported = Presentation(str(output_path))
    except Exception as exc:
        raise ValueError("Generated PPTX could not be reopened.") from exc
    if not exported.slides:
        raise ValueError("Generated PPTX does not contain any slides.")
    final_slide_text = [
        shape.text
        for shape in _iter_text_shapes(exported.slides[-1].shapes)
        if shape.text.strip()
    ]
    if "근거 출처" not in final_slide_text:
        raise ValueError("Generated PPTX is missing the final evidence slide.")


def _truncate(value: str, max_length: int = 180) -> str:
    normalized = " ".join(value.split())
    if len(normalized) <= max_length:
        return normalized
    return normalized[: max_length - 3].rstrip() + "..."
