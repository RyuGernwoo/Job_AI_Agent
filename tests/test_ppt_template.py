import io
import os
import sys
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from fastapi.testclient import TestClient
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

from lectureops_agent.app.main import create_app
from lectureops_agent.models.schemas import MaterialChunk, NCSUnit, ProjectCreate
from lectureops_agent.services.ppt_template_service import (
    InMemoryPPTTemplateStore,
    SupabasePPTTemplateStore,
    analyze_ppt_template,
)


def sample_project() -> ProjectCreate:
    return ProjectCreate(
        course_type="ncs",
        course_title="Programming Basics",
        lesson_title="Using Python Functions",
        learner_profile="Beginning Python learners",
        learning_objectives=["Explain function inputs and return values."],
        ncs_units=[
            NCSUnit(
                unit_code="2001020231_23v5",
                unit_name="프로그래밍 언어 활용",
                elements=["스크립트 언어 활용하기"],
                target_criteria=["스크립트 언어의 기본 문법을 활용할 수 있다."],
            )
        ],
    )


def template_bytes() -> bytes:
    presentation = Presentation()
    sample_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    sample_slide.shapes.title.text = "REMOVE THIS SAMPLE SLIDE"
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def create_test_client(
    *,
    template_store: InMemoryPPTTemplateStore | None = None,
) -> TestClient:
    with patch.dict(
        os.environ,
        {"LESSONPACK_ENV_FILE": str(ROOT / "missing-test.env")},
        clear=True,
    ):
        return TestClient(create_app(ppt_template_store=template_store))


class FakeSupabaseClient:
    def __init__(self) -> None:
        self.rows: dict[str, list[dict]] = {}
        self.storage = FakeStorage()

    def table(self, table_name: str):
        return FakeTable(self, table_name)


class FakeTable:
    def __init__(self, client: FakeSupabaseClient, table_name: str) -> None:
        self.client = client
        self.table_name = table_name
        self.filters: dict[str, object] = {}
        self.pending_row: dict | None = None
        self.pending_update: dict | None = None
        self.delete_requested = False

    def select(self, columns: str):
        return self

    def eq(self, key: str, value: object):
        self.filters[key] = value
        return self

    def limit(self, count: int):
        self.limit_count = count
        return self

    def upsert(self, row: dict, *, on_conflict: str):
        self.pending_row = row
        self.on_conflict = on_conflict
        return self

    def update(self, row: dict):
        self.pending_update = row
        return self

    def delete(self):
        self.delete_requested = True
        return self

    def execute(self):
        rows = self.client.rows.setdefault(self.table_name, [])
        if self.pending_row is not None:
            conflict_value = self.pending_row[self.on_conflict]
            rows[:] = [row for row in rows if row.get(self.on_conflict) != conflict_value]
            rows.append(dict(self.pending_row))
            return SimpleNamespace(data=[self.pending_row])
        selected = [
            row for row in rows if all(row.get(key) == value for key, value in self.filters.items())
        ]
        if self.pending_update is not None:
            for row in selected:
                row.update(self.pending_update)
            return SimpleNamespace(data=selected)
        if self.delete_requested:
            rows[:] = [row for row in rows if row not in selected]
            return SimpleNamespace(data=selected)
        return SimpleNamespace(data=selected[: getattr(self, "limit_count", len(selected))])


class FakeStorage:
    def __init__(self) -> None:
        self.buckets: dict[str, dict[str, bytes]] = {}

    def from_(self, bucket_name: str):
        return FakeBucket(self.buckets.setdefault(bucket_name, {}))


class FakeBucket:
    def __init__(self, objects: dict[str, bytes]) -> None:
        self.objects = objects

    def upload(self, path: str, content: bytes, file_options=None):
        if path in self.objects:
            raise RuntimeError("object already exists")
        self.objects[path] = bytes(content)
        return {"path": path}

    def download(self, path: str) -> bytes:
        return self.objects[path]

    def remove(self, paths: list[str]):
        for path in paths:
            self.objects.pop(path, None)
        return []


class PPTTemplateTests(unittest.TestCase):
    def test_template_analysis_builds_layout_manifest_and_mapping(self):
        metadata = analyze_ppt_template(
            project_id="project-001",
            filename="기관 강의 템플릿.pptx",
            content=template_bytes(),
        )

        self.assertEqual(metadata.project_id, "project-001")
        self.assertEqual(metadata.original_filename, "기관 강의 템플릿.pptx")
        self.assertEqual(metadata.source_slide_count, 1)
        self.assertGreater(len(metadata.layouts), 0)
        self.assertEqual(
            set(metadata.layout_mapping),
            {
                "cover",
                "objectives",
                "lesson",
                "practice",
                "assessment",
                "ncs_coverage",
                "sources",
            },
        )
        self.assertTrue(metadata.warnings)

    def test_template_analysis_rejects_non_pptx_content(self):
        with self.assertRaisesRegex(ValueError, "Office Open XML"):
            analyze_ppt_template(
                project_id="project-001",
                filename="invalid.pptx",
                content=b"not-a-pptx",
            )

    def test_template_api_upload_mapping_export_and_delete(self):
        store = InMemoryPPTTemplateStore()
        client = create_test_client(template_store=store)
        created = client.post("/api/projects", json=sample_project().model_dump(mode="json"))
        project_id = created.json()["project_id"]

        uploaded = client.post(
            f"/api/projects/{project_id}/ppt-template",
            files={
                "file": (
                    "기관 템플릿.pptx",
                    template_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

        self.assertEqual(uploaded.status_code, 200)
        metadata = uploaded.json()
        self.assertEqual(metadata["project_id"], project_id)
        self.assertEqual(metadata["status"], "ready")
        self.assertEqual(metadata["source_slide_count"], 1)

        body_layout = next(
            layout["layout_index"]
            for layout in metadata["layouts"]
            if layout["supports_title"] and layout["supports_body"]
        )
        blank_layout = next(
            layout["layout_index"]
            for layout in metadata["layouts"]
            if not layout["supports_title"] and not layout["supports_body"]
        )
        updated = client.put(
            f"/api/projects/{project_id}/ppt-template/mapping",
            json={"layout_mapping": {"lesson": body_layout, "sources": blank_layout}},
        )
        self.assertEqual(updated.status_code, 200)
        self.assertEqual(updated.json()["layout_mapping"]["lesson"], body_layout)

        chunk = MaterialChunk(
            chunk_id="doc001-p000-c001",
            project_id=project_id,
            document_id="doc001",
            source_name="sample.md",
            source_type="md",
            text="A function receives input and returns output.",
            metadata={"license": "PSF License"},
        )
        generated = client.post(
            f"/api/projects/{project_id}/generate",
            json={"retrieved_chunks": [chunk.model_dump(mode="json")]},
        )
        package_id = generated.json()["package_id"]
        exported = client.get(f"/api/packages/{package_id}/export.pptx")

        self.assertEqual(exported.status_code, 200)
        self.assertEqual(exported.headers["x-lessonpack-ppt-template-mode"], "custom")
        presentation = Presentation(io.BytesIO(exported.content))
        slide_text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertIn("Using Python Functions", slide_text)
        self.assertNotIn("REMOVE THIS SAMPLE SLIDE", slide_text)
        final_slide_text = [
            shape.text
            for shape in presentation.slides[-1].shapes
            if hasattr(shape, "text")
        ]
        self.assertIn("근거 출처", final_slide_text)

        deleted = client.delete(f"/api/projects/{project_id}/ppt-template")
        self.assertEqual(deleted.status_code, 204)
        missing = client.get(f"/api/projects/{project_id}/ppt-template")
        self.assertEqual(missing.status_code, 404)

    def test_template_mapping_rejects_unknown_layout_index(self):
        client = create_test_client()
        created = client.post("/api/projects", json=sample_project().model_dump(mode="json"))
        project_id = created.json()["project_id"]
        client.post(
            f"/api/projects/{project_id}/ppt-template",
            files={"file": ("template.pptx", template_bytes())},
        )

        response = client.put(
            f"/api/projects/{project_id}/ppt-template/mapping",
            json={"layout_mapping": {"lesson": 9999}},
        )

        self.assertEqual(response.status_code, 422)
        self.assertIn("do not exist", response.json()["detail"])

    def test_template_export_falls_back_to_default_when_content_load_fails(self):
        class BrokenContentStore(InMemoryPPTTemplateStore):
            def load_content(self, project_id: str) -> bytes:
                raise RuntimeError("storage unavailable")

        store = BrokenContentStore()
        client = create_test_client(template_store=store)
        created = client.post("/api/projects", json=sample_project().model_dump(mode="json"))
        project_id = created.json()["project_id"]
        client.post(
            f"/api/projects/{project_id}/ppt-template",
            files={"file": ("template.pptx", template_bytes())},
        )
        chunk = MaterialChunk(
            chunk_id="doc001-p000-c001",
            project_id=project_id,
            document_id="doc001",
            source_name="sample.md",
            source_type="md",
            text="A function receives input and returns output.",
        )
        generated = client.post(
            f"/api/projects/{project_id}/generate",
            json={"retrieved_chunks": [chunk.model_dump(mode="json")]},
        )

        exported = client.get(
            f"/api/packages/{generated.json()['package_id']}/export.pptx"
        )

        self.assertEqual(exported.status_code, 200)
        self.assertEqual(
            exported.headers["x-lessonpack-ppt-template-mode"],
            "default-fallback",
        )
        self.assertTrue(exported.content.startswith(b"PK"))

    def test_supabase_template_store_round_trip(self):
        client = FakeSupabaseClient()
        store = SupabasePPTTemplateStore(client=client)
        content = template_bytes()
        metadata = analyze_ppt_template(
            project_id="project-001",
            filename="template.pptx",
            content=content,
        )

        saved = store.save(metadata, content)
        loaded = store.get("project-001")
        loaded_content = store.load_content("project-001")
        updated = store.update_mapping(
            "project-001",
            {**metadata.layout_mapping, "lesson": metadata.layouts[0].layout_index},
        )

        self.assertEqual(saved, metadata)
        self.assertEqual(loaded, metadata)
        self.assertEqual(loaded_content, content)
        self.assertIsNotNone(updated)
        self.assertEqual(updated.layout_mapping["lesson"], metadata.layouts[0].layout_index)
        self.assertTrue(store.delete("project-001"))
        self.assertIsNone(store.get("project-001"))


if __name__ == "__main__":
    unittest.main()
