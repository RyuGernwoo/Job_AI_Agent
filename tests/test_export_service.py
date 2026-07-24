import io
import sys
import tempfile
import unittest
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

from lectureops_agent.models.schemas import MaterialChunk, NCSUnit, PackageStatus, ProjectCreate
from lectureops_agent.services.export_service import (
    build_export_filename,
    export_lesson_package_docx,
    export_lesson_package_pptx,
)
from lectureops_agent.services.generation_service import generate_lesson_package
from lectureops_agent.services.ppt_template_service import analyze_ppt_template


def make_package(status: PackageStatus = PackageStatus.GENERATED):
    project = ProjectCreate(
        course_type="ncs",
        course_title="Generative AI Python Basics",
        lesson_title="Python functions and prompt automation practice",
        learner_profile="Job training learners with basic Python experience",
        total_training_hours=8,
        total_lessons=4,
        theory_ratio_percent=30,
        practice_ratio_percent=70,
        learning_objectives=["Explain function inputs and return values."],
        ncs_units=[
            NCSUnit(
                unit_code="MVP-NCS-001",
                unit_name="AI basics",
                elements=["Explain basic AI concepts."],
            )
        ],
    ).to_project(project_id="project-001")
    chunk = MaterialChunk(
        chunk_id="doc001-p000-c001",
        project_id=project.project_id,
        document_id="doc001",
        source_name="sample.pdf",
        source_type="pdf",
        page=None,
        text="A function can receive input and return output.",
        metadata={
            "license": "PSF License",
            "source_url": "https://docs.python.org/3/tutorial/controlflow.html",
            "source_file": "data/raw/materials/tutorial_functions.md",
        },
    )
    package = generate_lesson_package(project=project, retrieved_chunks=[chunk])
    return package.model_copy(update={"status": status})


class ExportServiceTests(unittest.TestCase):
    def test_export_lesson_package_docx_writes_readable_document(self):
        package = make_package(PackageStatus.GENERATED)
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "lesson_package.docx"

            result = export_lesson_package_docx(package=package, output_path=output_path)

            self.assertEqual(result, output_path)
            self.assertTrue(output_path.exists())
            doc = Document(str(output_path))
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            self.assertIn("Python functions and prompt automation practice", text)
            self.assertIn("실습 과제", text)
            self.assertIn("훈련 운영: 총 8시간 · 4차시 · 이론 30% · 실습 70%", text)
            self.assertIn("근거 출처", text)
            self.assertIn("PSF License", text)
            self.assertIn("NCS 연계", text)
            self.assertNotIn("근거:", text)
            self.assertNotIn("검수 이력", text)
            self.assertNotIn(package.package_id, text)
            self.assertEqual(text.count("근거 출처"), 1)

    def test_export_lesson_package_pptx_writes_summary_slides(self):
        package = make_package(PackageStatus.REGENERATED)
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "lesson_package.pptx"

            result = export_lesson_package_pptx(package=package, output_path=output_path)

            self.assertEqual(result, output_path)
            self.assertTrue(output_path.exists())
            presentation = Presentation(str(output_path))
            slide_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("Python functions and prompt automation practice", slide_text)
            self.assertIn("실습 개요", slide_text)
            self.assertIn("실습 수행 절차", slide_text)
            # 객관식 보기(선택지)와 정답이 슬라이드에 포함되는지 확인한다.
            self.assertIn("정답: 1번", slide_text)
            self.assertIn("훈련 운영: 총 8시간 · 4차시 · 이론 30% · 실습 70%", slide_text)
            self.assertIn("근거 출처", slide_text)
            self.assertIn("PSF License", slide_text)
            self.assertNotIn("근거:", slide_text)
            self.assertNotIn("검수 이력", slide_text)
            self.assertNotIn(package.package_id, slide_text)
            self.assertEqual(presentation.slides[-1].shapes.title.text, "근거 출처")
            for slide in list(presentation.slides)[:-1]:
                text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                self.assertNotIn(chunk_id := package.evidence_sources[0].chunk_id, text)
            self.assertIn(chunk_id, "\n".join(shape.text for shape in presentation.slides[-1].shapes if hasattr(shape, "text")))

    def test_export_uses_distinct_source_slide_designs_beyond_cover(self):
        template = Presentation()
        slide_specs = [
            ("Course Cover", "Course subtitle", "COVER_ACCENT"),
            ("Contents", "Learning objectives", "OBJECTIVES_ACCENT"),
            ("Long Content", "Lesson explanation", "LESSON_ACCENT"),
            ("Timeline STEP1", "Practice procedure", "PRACTICE_ACCENT"),
            ("Quick Check Quiz", "Assessment question", "ASSESSMENT_ACCENT"),
            ("Core Skills", "NCS performance criteria", "NCS_ACCENT"),
            ("References and Sources", "Source list", "SOURCES_ACCENT"),
        ]
        for title, body, accent_name in slide_specs:
            slide = template.slides.add_slide(template.slide_layouts[6])
            slide.shapes.add_textbox(
                Inches(0.8),
                Inches(0.6),
                Inches(8),
                Inches(0.8),
            ).text_frame.text = title
            slide.shapes.add_textbox(
                Inches(0.8),
                Inches(1.8),
                Inches(8),
                Inches(3.8),
            ).text_frame.text = body
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(9),
                Inches(0),
                Inches(1),
                Inches(7.5),
            )
            accent.name = accent_name
        stream = io.BytesIO()
        template.save(stream)
        content = stream.getvalue()
        metadata = analyze_ppt_template(
            project_id="project-source-layouts",
            filename="source-layouts.pptx",
            content=content,
        )
        package = make_package(PackageStatus.GENERATED)

        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "source_layouts.pptx"
            export_lesson_package_pptx(
                package=package,
                output_path=output_path,
                template_content=content,
                layout_mapping=metadata.layout_mapping,
            )
            exported = Presentation(str(output_path))

        shape_names = {
            shape.name
            for slide in exported.slides
            for shape in slide.shapes
        }
        self.assertTrue(
            {
                "COVER_ACCENT",
                "OBJECTIVES_ACCENT",
                "LESSON_ACCENT",
                "PRACTICE_ACCENT",
                "ASSESSMENT_ACCENT",
                "NCS_ACCENT",
                "SOURCES_ACCENT",
            }.issubset(shape_names)
        )
        slide_text = "\n".join(
            shape.text
            for slide in exported.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertNotIn("Lesson explanation", slide_text)
        self.assertNotIn("Assessment question", slide_text)

    def test_build_export_filename_uses_safe_lesson_title(self):
        package = make_package(PackageStatus.GENERATED)
        package.lesson_plan.title = 'Python 함수/자료구조: 자동화 실습? "입문"'

        docx_name = build_export_filename(package, ".docx")
        pptx_name = build_export_filename(package, "pptx")

        self.assertEqual(docx_name, "Python_함수_자료구조_자동화_실습_입문_교안.docx")
        self.assertEqual(pptx_name, "Python_함수_자료구조_자동화_실습_입문_강의자료.pptx")
        self.assertNotIn(package.package_id, docx_name)

if __name__ == "__main__":
    unittest.main()
