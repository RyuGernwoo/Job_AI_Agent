from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation


ENGLISH_LABELS = [
    "Learning Objectives",
    "Lesson Plan",
    "Practice",
    "Assessment",
    "Submission",
    "Answer",
    "Explanation",
    "Citations",
    "Evidence Sources",
]
UUID_PATTERN = re.compile(r"\b[0-9a-f]{8}-[0-9a-f]{4}-[1-5][0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}\b", re.IGNORECASE)
CITATION_PATTERN = re.compile(r"\b[0-9A-Za-z_-]+(?:-p\d+)?-c\d+\b")
DUPLICATE_LABELS = ["실습 실습", "제출물: 제출물:", "수행 절차: 수행 절차:", "평가 기준: 평가 기준:"]
AWKWARD_PHRASES = ["있다. 및", "한다. 및"]


def inspect_exports(*, docx_path: Path | None, pptx_path: Path | None) -> dict[str, Any]:
    report: dict[str, Any] = {"passed": True, "files": {}}
    if docx_path:
        report["files"]["docx"] = _inspect_text(_docx_text(docx_path))
    if pptx_path:
        report["files"]["pptx"] = _inspect_text(_pptx_text(pptx_path))

    for result in report["files"].values():
        if not result["passed"]:
            report["passed"] = False
    return report


def _docx_text(path: Path) -> str:
    document = Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _pptx_text(path: Path) -> str:
    presentation = Presentation(path)
    values: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                values.append(shape.text)
    return "\n".join(values)


def _inspect_text(text: str) -> dict[str, Any]:
    evidence_index = text.rfind("근거 출처")
    content_before_evidence = text[:evidence_index] if evidence_index >= 0 else text
    signals = {
        "has_korean_objectives": "학습 목표" in text,
        "has_evidence_sources": "근거 출처" in text,
        "has_license": "라이선스:" in text,
        "has_ncs_alignment": "NCS 연계:" in text,
        "has_review_section": "검수 이력" in text,
        "has_english_labels": any(label in text for label in ENGLISH_LABELS),
        "has_inline_citation_labels": "근거:" in content_before_evidence,
        "has_citation_ids_before_evidence": bool(CITATION_PATTERN.search(content_before_evidence)),
        "has_internal_uuid": bool(UUID_PATTERN.search(text)),
        "has_duplicate_labels": any(label in text for label in DUPLICATE_LABELS),
        "has_awkward_phrases": any(phrase in text for phrase in AWKWARD_PHRASES),
        "evidence_section_count": text.count("근거 출처"),
    }
    passed = (
        signals["has_korean_objectives"]
        and signals["has_evidence_sources"]
        and signals["has_ncs_alignment"]
        and not signals["has_english_labels"]
        and not signals["has_inline_citation_labels"]
        and not signals["has_citation_ids_before_evidence"]
        and not signals["has_internal_uuid"]
        and not signals["has_duplicate_labels"]
        and not signals["has_awkward_phrases"]
        and signals["evidence_section_count"] == 1
    )
    return {"passed": passed, "signals": signals}


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect LessonPack AI DOCX/PPTX export quality.")
    parser.add_argument("--docx", type=Path)
    parser.add_argument("--pptx", type=Path)
    parser.add_argument("--report", type=Path)
    args = parser.parse_args()

    if not args.docx and not args.pptx:
        parser.error("at least one of --docx or --pptx is required")

    report = inspect_exports(docx_path=args.docx, pptx_path=args.pptx)
    text = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(text + "\n", encoding="utf-8")
    print(text)
    return 0 if report["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
